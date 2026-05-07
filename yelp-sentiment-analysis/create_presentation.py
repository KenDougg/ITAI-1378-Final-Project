#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Yelp Sentiment Analysis project.
Requires: pip install python-pptx

Run: python create_presentation.py
Output: docs/Yelp_Sentiment_Analysis_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(0, 51, 102)      # Dark blue
LIGHT_BLUE = RGBColor(0, 102, 204)    # Light blue
RED = RGBColor(214, 39, 40)           # Red for negative
ORANGE = RGBColor(255, 127, 14)       # Orange for neutral
GREEN = RGBColor(44, 160, 44)         # Green for positive
WHITE = RGBColor(255, 255, 255)       # White
DARK_GRAY = RGBColor(64, 64, 64)      # Dark gray
LIGHT_GRAY = RGBColor(240, 240, 240)  # Light gray


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True

    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, content_items, bg_color=WHITE):
    """Add a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE

    # Add title text
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title_frame.margin_left = Inches(0.5)
    title_frame.margin_top = Inches(0.2)

    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.7))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()

        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)

    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE

    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.margin_left = Inches(0.5)
    title_frame.margin_top = Inches(0.15)

    # Left column - title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4.5), Inches(0.5))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = left_title
    left_title_frame.paragraphs[0].font.size = Pt(22)
    left_title_frame.paragraphs[0].font.bold = True
    left_title_frame.paragraphs[0].font.color.rgb = GREEN

    # Left column - content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.5), Inches(5.4))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for i, item in enumerate(left_items):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    # Right column - title
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.0), Inches(4.3), Inches(0.5))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = right_title
    right_title_frame.paragraphs[0].font.size = Pt(22)
    right_title_frame.paragraphs[0].font.bold = True
    right_title_frame.paragraphs[0].font.color.rgb = RED

    # Right column - content
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(5.4))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    for i, item in enumerate(right_items):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    return slide


# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
add_title_slide(
    prs,
    "Sentiment Analysis",
    "Analyzing Customer Reviews from Yelp\n\nGroup 5: Esbeide Ibarra, Joel Garza, Nhat Nam Khanh Nguyen\nProfessor: Viswanatha Rao"
)

# ============================================================================
# SLIDE 2: Problem Definition
# ============================================================================
add_content_slide(
    prs,
    "The Problem",
    [
        "🔴 Challenge: Businesses receive thousands of customer reviews monthly",
        "",
        "❓ Question: How can we automatically understand customer sentiment?",
        "",
        "⏱️  Issue: Manual review analysis is time-consuming (~1 hour per 100 reviews)",
        "",
        "💡 Opportunity: Build an ML system to classify sentiment automatically",
        "",
        "💰 Impact: Save 20+ hours/month + identify key satisfaction drivers"
    ]
)

# ============================================================================
# SLIDE 3: Business Impact
# ============================================================================
add_content_slide(
    prs,
    "Why Does This Matter?",
    [
        "✅ Triage Negative Reviews: Auto-flag 1-2 star reviews for urgent response",
        "",
        "📊 Track Trends: Monitor monthly sentiment score to measure improvement",
        "",
        "🎯 Target Marketing: Extract positive keywords for campaigns",
        "",
        "🔍 Root Cause Analysis: Identify common complaints across restaurants",
        "",
        "⚡ Real-time Response: Respond to critical feedback within hours, not days"
    ]
)

# ============================================================================
# SLIDE 4: Dataset Overview
# ============================================================================
add_content_slide(
    prs,
    "Our Dataset: Yelp Reviews",
    [
        "📊 Source: Hugging Face yelp_review_full (650,000+ reviews)",
        "",
        "📝 Sample Used: 1,000 reviews for manageable training",
        "",
        "🎯 Classes:",
        "    • Negative (1-2 stars): 404 reviews - 40.4%",
        "    • Positive (4-5 stars): 392 reviews - 39.2%",
        "    • Neutral (3 stars): 204 reviews - 20.4%",
        "",
        "📈 Train/Test Split: 80% (800) / 20% (200)"
    ]
)

# ============================================================================
# SLIDE 5: Technical Approach
# ============================================================================
add_content_slide(
    prs,
    "Our Approach: TF-IDF + Machine Learning",
    [
        "1️⃣  Preprocessing: Lowercase, remove URLs, tokenize, lemmatize",
        "",
        "2️⃣  Feature Extraction: Convert text to 5,000 TF-IDF features",
        "",
        "3️⃣  Model Training: Train two models for comparison",
        "    • Logistic Regression",
        "    • Multinomial Naive Bayes",
        "",
        "4️⃣  Evaluation: Test on unseen data with multiple metrics",
        "",
        "5️⃣  Prediction: Classify new reviews as Positive/Negative/Neutral"
    ]
)

# ============================================================================
# SLIDE 6: Models Trained
# ============================================================================
add_two_column_slide(
    prs,
    "Two Models Trained & Compared",
    "Logistic Regression",
    [
        "✓ Interpretable",
        "✓ Fast training",
        "Accuracy: 67.5%",
        "Precision: 54.2%",
        "Recall: 67.5%",
        "F1-Score: 60.1%"
    ],
    "Naive Bayes (Best)",
    [
        "✓ Very fast (0.02 sec)",
        "✓ Competitive accuracy",
        "Accuracy: 68.0% 🏆",
        "Precision: 54.1%",
        "Recall: 68.0%",
        "F1-Score: 60.2%"
    ]
)

# ============================================================================
# SLIDE 7: Results - What Worked
# ============================================================================
add_content_slide(
    prs,
    "What Worked Well ✓",
    [
        "🟢 Positive Reviews: 86% correctly identified",
        "    → Model easily recognizes 'great', 'excellent', 'amazing'",
        "",
        "🔴 Negative Reviews: 86% correctly identified",
        "    → Model easily recognizes 'terrible', 'bad', 'disappointing'",
        "",
        "⚡ Fast Training: Naive Bayes trained in 0.02 seconds",
        "",
        "📊 Clear Separation: Strong vs. weak sentiments naturally separable",
        "",
        "✅ Overall Accuracy: 68% on completely unseen test data"
    ]
)

# ============================================================================
# SLIDE 8: Results - What Failed
# ============================================================================
add_content_slide(
    prs,
    "The Challenge: Neutral Class ✗",
    [
        "🟡 Neutral Reviews (3 stars): 0% accuracy (all 41 samples misclassified)",
        "",
        "❌ Why?",
        "    1. Class Imbalance: Only 204 neutral vs 400+ positive/negative",
        "    2. Ambiguous Language: '3-star' reviews mix good & bad",
        "    3. TF-IDF Limitation: Can't understand context ('good food, bad service')",
        "",
        "📌 Root Cause: When training data has fewer examples, model defaults to common classes",
        "",
        "💡 Silver Lining: This is a real, documented problem—not a bug, it's insight!"
    ]
)

# ============================================================================
# SLIDE 9: Success & Failure Examples
# ============================================================================
add_two_column_slide(
    prs,
    "Real Examples: Successes vs. Failures",
    "✅ SUCCESS",
    [
        'Text: "Love this place! Best restaurant in town."',
        "Predicted: POSITIVE ✓",
        "Confidence: 92%",
        "",
        'Text: "Terrible experience. Cold food, rude staff."',
        "Predicted: NEGATIVE ✓",
        "Confidence: 87%"
    ],
    "❌ FAILURE",
    [
        'Text: "Nice place, but food quality is mediocre."',
        "Actual: NEUTRAL",
        "Predicted: POSITIVE ✗",
        "Confidence: 72%",
        "",
        "[All 41 neutral reviews misclassified]"
    ]
)

# ============================================================================
# SLIDE 10: Key Learnings
# ============================================================================
add_content_slide(
    prs,
    "Key Learnings",
    [
        "🧠 Text Preprocessing Matters: Lemmatization + stopwords = big accuracy boost",
        "",
        "🎯 Simple Models Work: Naive Bayes competitive with complex methods",
        "",
        "⚖️  Class Imbalance is Real: Minority classes get ignored by default",
        "",
        "📚 TF-IDF Limitations: Ignores word order, context, sarcasm",
        "",
        "✍️  Honest Analysis > Perfect Metrics: Understanding failures is valuable"
    ]
)

# ============================================================================
# SLIDE 11: Future Improvements
# ============================================================================
add_content_slide(
    prs,
    "Next Steps: Future Improvements",
    [
        "🔧 Quick Wins:",
        "    → Use SMOTE to balance classes (could reach 70-72% accuracy)",
        "    → Implement class weights in model",
        "",
        "🚀 Advanced:",
        "    → BERT Transformers: Understand context & sarcasm (75%+ accuracy)",
        "    → Aspect-Based Sentiment: Separate scores for food/service/ambiance",
        "    → Collect more neutral examples: Better train the minority class",
        "",
        "📈 Production:",
        "    → Deploy with confidence thresholds",
        "    → Build API for business integration"
    ]
)

# ============================================================================
# SLIDE 12: Conclusion & Questions
# ============================================================================
add_title_slide(
    prs,
    "Questions?",
    "68% Accuracy • Strong Foundation • Clear Path Forward\n\nRepository: github.com/[group5]/yelp-sentiment-analysis\nThanks to Professor Viswanatha Rao"
)

# ============================================================================
# SAVE PRESENTATION
# ============================================================================
from pathlib import Path
output_dir = Path(__file__).parent / "docs"
output_dir.mkdir(parents=True, exist_ok=True)
output_path = output_dir / "Yelp_Sentiment_Analysis_Presentation.pptx"
prs.save(str(output_path))
print(f"✅ Presentation created successfully!")
print(f"📍 Location: {output_path}")
print(f"📊 Slides: 12")
print(f"⏱️  Estimated time: 10 minutes")
print(f"\nTo view: Open {output_path} in PowerPoint or Google Slides")
